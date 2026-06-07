#!/usr/bin/env python3
"""Inject markdown speaker notes into PPTX speaker notes pane."""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Install: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

SLIDE_HEADER = re.compile(r"^## Slide (\d+)", re.MULTILINE)


def parse_notes(md_path: Path) -> dict[int, str]:
    content = md_path.read_text()
    # Stop at appendix
    if "## Appendix" in content:
        content = content.split("## Appendix")[0]

    notes: dict[int, str] = {}
    parts = re.split(r"\n## Slide (\d+)", content)
    for i in range(1, len(parts), 2):
        num = int(parts[i])
        body = parts[i + 1]
        body = body.split("\n---\n")[0].strip()
        # Drop title line (first line after header often starts with —)
        lines = body.split("\n")
        cleaned: list[str] = []
        for line in lines:
            if re.match(r"^#\s", line):
                continue
            if re.match(r"^\*\*~\d+", line):
                continue
            cleaned.append(line)
        text = "\n".join(cleaned).strip()
        if text:
            notes[num] = text
    return notes


def inject(pptx_path: Path, notes: dict[int, str], output: Path) -> None:
    prs = Presentation(str(pptx_path))
    for i, slide in enumerate(prs.slides, 1):
        note = notes.get(i)
        if not note:
            continue
        tf = slide.notes_slide.notes_text_frame
        tf.clear()
        chunks = [note[j : j + 3000] for j in range(0, len(note), 3000)]
        for ci, chunk in enumerate(chunks):
            p = tf.paragraphs[0] if ci == 0 else tf.add_paragraph()
            p.text = chunk
    prs.save(str(output))
    print(f"Wrote {len(notes)} slide notes to {output}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Inject markdown notes into PPTX")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("markdown", type=Path, help="Speaker notes markdown")
    parser.add_argument(
        "--output", "-o", type=Path, default=None, help="Default: overwrite pptx"
    )
    args = parser.parse_args()

    if not args.pptx.is_file():
        print(f"Not found: {args.pptx}", file=sys.stderr)
        sys.exit(1)
    if not args.markdown.is_file():
        print(f"Not found: {args.markdown}", file=sys.stderr)
        sys.exit(1)

    notes = parse_notes(args.markdown)
    out = args.output or args.pptx
    inject(args.pptx, notes, out)


if __name__ == "__main__":
    main()
