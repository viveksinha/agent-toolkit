#!/usr/bin/env python3
"""Extract text, existing notes, and images from a PPTX deck."""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Install: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def shape_text(shape) -> list[str]:
    texts: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append(t)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(shape_text(child))
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shape.table.rows:
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    texts.append(t)
    return texts


def save_pictures(slide, slide_num: int, images_dir: Path) -> list[str]:
    saved: list[str] = []
    idx = 0

    def walk(shapes):
        nonlocal idx
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ext = shape.image.content_type.split("/")[-1]
                path = images_dir / f"slide{slide_num:02d}_img{idx}.{ext}"
                path.write_bytes(shape.image.blob)
                saved.append(str(path))
                idx += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)

    walk(slide.shapes)
    return saved


def extract(pptx_path: Path, output: Path, images_dir: Path | None) -> int:
    prs = Presentation(str(pptx_path))
    lines: list[str] = []
    image_manifest: list[str] = []

    if images_dir:
        images_dir.mkdir(parents=True, exist_ok=True)

    for i, slide in enumerate(prs.slides, 1):
        lines.append("")
        lines.append("=" * 80)
        lines.append(f"SLIDE {i}")
        lines.append("=" * 80)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"[EXISTING NOTES]: {notes}")

        seen: set[str] = set()
        for shape in slide.shapes:
            for t in shape_text(shape):
                if t not in seen:
                    seen.add(t)
                    lines.append(t)

        if images_dir:
            imgs = save_pictures(slide, i, images_dir)
            if imgs:
                lines.append(f"[IMAGES]: {', '.join(imgs)}")
                image_manifest.extend(imgs)

        if len(seen) <= 3:
            lines.append("[LOW TEXT — review images if exported]")

    output.write_text("\n".join(lines).lstrip() + "\n")
    print(f"Slides: {len(prs.slides)}")
    print(f"Text: {output}")
    if images_dir:
        print(f"Images: {len(image_manifest)} files in {images_dir}")
    return len(prs.slides)


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract PPTX slide content")
    parser.add_argument("pptx", type=Path, help="Path to .pptx file")
    parser.add_argument(
        "--output", "-o", type=Path, default=Path("/tmp/slide-extract.txt")
    )
    parser.add_argument(
        "--images-dir",
        type=Path,
        default=None,
        help="Export slide images for visual/diagram slides",
    )
    args = parser.parse_args()

    if not args.pptx.is_file():
        print(f"Not found: {args.pptx}", file=sys.stderr)
        sys.exit(1)

    extract(args.pptx, args.output, args.images_dir)


if __name__ == "__main__":
    main()
